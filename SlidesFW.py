import openai
import json
from pptx import Presentation  # You can install it using 'pip install python-pptx'

# Set your OpenAI API key
openai.api_key = "API KEY"

# Ask the user for the presentation title in Terminal
presentation_title = input("What do you want to make a presentation about?: ")

# Define the JSON structure for the OpenAI request
query_json = """{
    "input_text": "[[QUERY]]",
    "output_format": "json",
    "json_structure": {
        "slides":"{{presentation_slides}}"
       }
    }"""

# Define the question prompt for generating the presentation
question = (
    "Generate a 10 slide presentation for the topic. Produce 100 to 200 words per slide. "
    + presentation_title
    + ". Each slide should have a {{header}}, {{content}}. The header is a subject from the topic and the content is interesting information about the header. The final slide should be a list of discussion questions. Return as JSON."
)

# Replace the placeholder in the JSON structure with the question
prompt = query_json.replace("[[QUERY]]", question)

# Display the prompt
print(prompt)

# Use OpenAI's ChatCompletion API to generate the presentation content
completion = openai.ChatCompletion.create(
    model="gpt-3.5-turbo", messages=[{"role": "user", "content": prompt}]
)

# Retrieve the response from the API
response = completion.choices[0].message.content

# Print the response
print(response)

# Parse the JSON response
r = json.loads(response)

# Extract the slide data
slide_data = r["slides"]

# Create a PowerPoint presentation
prs = Presentation()

# Iterate through the slide data and add slides to the presentation
for slide in slide_data:
    slide_layout = prs.slide_layouts[1]
    new_slide = prs.slides.add_slide(slide_layout)

    # Add header if available
    if slide["header"]:
        title = new_slide.shapes.title
        title.text = slide["header"]

    # Add content if available
    if slide["content"]:
        shapes = new_slide.shapes
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = slide["content"]
        tf.fit_text(font_family="Calibri", max_size=18, bold=True)

# Save the PowerPoint presentation as 'output.pptx'
prs.save("output.pptx")
